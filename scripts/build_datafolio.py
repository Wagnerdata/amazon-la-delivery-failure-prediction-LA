"""
build_datafolio.py — Generate deliverables/05_datafolio.pptx

Single-slide conference poster (24" × 36" portrait) for the
Correlation One Data Analytics portfolio project.

Charts must already exist in reports/figures/:
  class_imbalance.png, correlation_heatmap.png,
  failure_rate_by_carrier.png, failure_rate_by_shift.png,
  failure_rate_by_distance.png, feature_importance_preview.png
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

# ── Paths ──────────────────────────────────────────────────────────────────────
BASE         = Path(__file__).parent.parent
FIGURES_DIR  = BASE / "reports" / "figures"
OUT_PATH     = BASE / "deliverables" / "05_datafolio.pptx"
OUT_PATH.parent.mkdir(parents=True, exist_ok=True)

# ── Color palette ──────────────────────────────────────────────────────────────
C_NAVY    = RGBColor(0x2C, 0x3E, 0x50)
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_DANGER  = RGBColor(0xE7, 0x4C, 0x3C)
C_WARN    = RGBColor(0xF3, 0x9C, 0x12)
C_SAFE    = RGBColor(0x27, 0xAE, 0x60)
C_LIGHT   = RGBColor(0xEC, 0xF0, 0xF1)
C_ACCENT  = RGBColor(0x29, 0x80, 0xB9)

# ── Slide dimensions: 24" × 36" portrait ──────────────────────────────────────
W_IN  = 24.0
H_IN  = 36.0

prs = Presentation()
prs.slide_width  = Inches(W_IN)
prs.slide_height = Inches(H_IN)

slide_layout = prs.slide_layouts[6]   # blank
slide = prs.slides.add_slide(slide_layout)


# ── Helper: add a filled rectangle ────────────────────────────────────────────
def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(
        1,   # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    return shape


# ── Helper: add a text box ─────────────────────────────────────────────────────
def add_text(slide, text, left, top, width, height,
             font_size=10, bold=False, color=None, align=PP_ALIGN.LEFT,
             wrap=True):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap

    # Handle multi-line: split on \n and add paragraphs
    lines = text.split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = color
    return txBox


# ── Helper: add image ──────────────────────────────────────────────────────────
def add_image(slide, img_path, left, top, width, height):
    if not img_path.exists():
        print(f"  [WARNING] Image not found: {img_path.name}")
        return
    slide.shapes.add_picture(
        str(img_path),
        Inches(left), Inches(top), Inches(width), Inches(height),
    )


# ══════════════════════════════════════════════════════════════════════════════
# TOP BANNER (0 to 4.2 inches tall)
# ══════════════════════════════════════════════════════════════════════════════
add_rect(slide, 0, 0, W_IN, 4.2, C_NAVY)

# Main title
add_text(
    slide,
    "Predicting Last-Mile Delivery Failure Before the Truck Leaves",
    0.3, 0.25, 23.4, 1.5,
    font_size=34, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER,
)

# Subtitle
add_text(
    slide,
    "Evidence from 3,559 Real Amazon Packages — Los Angeles, July 2018",
    0.3, 1.55, 23.4, 0.8,
    font_size=20, bold=False, color=RGBColor(0xAE, 0xC6, 0xCF),
    align=PP_ALIGN.CENTER,
)

# Program line
add_text(
    slide,
    "Correlation One Data Analytics (DANA) — Week 12 Final Portfolio Project  |  Amazon Last Mile Routing Research Challenge Dataset",
    0.3, 2.25, 23.4, 0.6,
    font_size=12, bold=False, color=RGBColor(0xBD, 0xC3, 0xC7),
    align=PP_ALIGN.CENTER,
)

# Key stat boxes
stat_w = 4.6
stat_h = 1.2
stats = [
    ("3,559", "Total Packages"),
    ("15", "LA Routes"),
    ("0.70%", "Failure Rate"),
    ("~140:1", "Class Imbalance"),
    ("$17", "Cost per Failure"),
]
for i, (num, label) in enumerate(stats):
    x = 0.3 + i * (stat_w + 0.1)
    add_rect(slide, x, 2.85, stat_w, stat_h, C_ACCENT)
    add_text(slide, num, x, 2.90, stat_w, 0.70, font_size=26, bold=True,
             color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, label, x, 3.55, stat_w, 0.45, font_size=12, bold=False,
             color=C_WHITE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# FOUR COLUMNS — start at y=4.4, height ~31.3 inches
# ══════════════════════════════════════════════════════════════════════════════
COL_Y     = 4.40
COL_H     = 31.2
COL_GAP   = 0.15
COL_W     = (W_IN - COL_GAP * 5) / 4   # ≈ 5.7"
col_x     = [COL_GAP + i * (COL_W + COL_GAP) for i in range(4)]

# ── Section header helper ──────────────────────────────────────────────────────
def section_header(slide, title, x, y, w, color=C_NAVY):
    add_rect(slide, x, y, w, 0.55, color)
    add_text(slide, title, x + 0.1, y + 0.05, w - 0.2, 0.45,
             font_size=14, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)


# ═══════════════════════════════════════════════════
# COLUMN 1 — BACKGROUND
# ═══════════════════════════════════════════════════
cx = col_x[0]
cy = COL_Y

add_rect(slide, cx, cy, COL_W, COL_H, C_LIGHT)

section_header(slide, "BACKGROUND", cx, cy + 0.15, COL_W)

body1 = """\
Business Problem
Every failed delivery attempt costs an estimated $17 when redelivery, \
customer service, and downstream Voice-of-Customer impact are combined. \
At scale, even a sub-1% failure rate represents material operational waste.

Dataset
Amazon Last Mile Routing Research Challenge (LMRC) 2021 — real operational data, \
not synthetic. 15 routes in Los Angeles, July 2018. 3,559 packages. \
Target: delivery_failure (1 = DELIVERY_ATTEMPTED, 0 = DELIVERED).

Key Dataset Notes
• days_in_fc: zero variance — all packages dispatched same day
• weather_risk: zero variance — July LA is universally "low"
• Only 2 shifts observed: morning and afternoon (no night routes)
• 91% of packages have cr_number_missing = 1 (normal for this dataset)

Approach
Random Forest classifier trained on dispatch-time features only. \
Class imbalance handled via class_weight='balanced' and SMOTE comparison. \
Threshold tuned for maximum recall (cost of missing a failure >> false alarm cost).

Tools & Stack
• Python 3.9+  •  scikit-learn  •  imbalanced-learn (SMOTE)
• pandas / numpy  •  matplotlib / seaborn
• SQLite (EDA queries)  •  CrewAI agents (dispatcher reports)

Why Recall Over Accuracy
A naive classifier that always predicts "Delivered" achieves 99.3% accuracy — \
and catches zero failures. Recall is the business-relevant metric: \
each missed failure costs $17; a false alarm costs a quick manual check."""

add_text(slide, body1, cx + 0.15, cy + 0.90, COL_W - 0.3, 14.0,
         font_size=9.5, color=C_NAVY)

# Features list
section_header(slide, "FEATURES USED", cx, cy + 15.5, COL_W, C_ACCENT)
features_text = """\
Carrier (A / B / C / D)
Shift (morning / afternoon)
Package type (standard / high_value)
Route distance (km) — bucketed
Packages in route (count)
Double scan flag (binary)
Locker issue flag (binary)
CR number missing flag (binary)

Excluded Features
damaged_on_arrival → IS the target (data leakage)
weather_risk → zero variance (excluded)
days_in_fc → zero variance (excluded)"""

add_text(slide, features_text, cx + 0.15, cy + 16.15, COL_W - 0.3, 6.5,
         font_size=9.5, color=C_NAVY)

# Imbalance note
section_header(slide, "CLASS IMBALANCE PROBLEM", cx, cy + 23.0, COL_W, C_DANGER)
imb_text = """\
25 failures in 3,559 packages = ~140:1 ratio.

Standard accuracy is meaningless here. A model \
predicting "delivered" every time scores 99.3% \
accuracy yet catches zero failures.

Solutions applied:
1. class_weight='balanced' in RandomForest
2. SMOTE oversampling comparison
3. Threshold optimization for recall"""
add_text(slide, imb_text, cx + 0.15, cy + 23.65, COL_W - 0.3, 5.5,
         font_size=9.5, color=C_NAVY)


# ═══════════════════════════════════════════════════
# COLUMN 2 — DATA & METHOD
# ═══════════════════════════════════════════════════
cx = col_x[1]
cy = COL_Y

add_rect(slide, cx, cy, COL_W, COL_H, C_LIGHT)
section_header(slide, "DATA & METHOD", cx, cy + 0.15, COL_W)

# Class imbalance chart
add_text(slide, "Class Distribution", cx + 0.15, cy + 0.90, COL_W - 0.3, 0.4,
         font_size=11, bold=True, color=C_NAVY)
add_image(slide, FIGURES_DIR / "class_imbalance.png",
          cx + 0.10, cy + 1.30, COL_W - 0.2, 4.2)

add_text(slide,
         "25 failures vs 3,534 successful deliveries. The ~140:1 imbalance "
         "drives all modeling decisions: metric choice, sampling strategy, "
         "and threshold tuning.",
         cx + 0.15, cy + 5.60, COL_W - 0.3, 1.2,
         font_size=9, color=C_NAVY)

# Correlation heatmap
add_text(slide, "Correlation Heatmap", cx + 0.15, cy + 6.90, COL_W - 0.3, 0.4,
         font_size=11, bold=True, color=C_NAVY)
add_image(slide, FIGURES_DIR / "correlation_heatmap.png",
          cx + 0.10, cy + 7.30, COL_W - 0.2, 5.2)

add_text(slide,
         "No single feature dominates in Pearson correlation — expected at 0.7% "
         "failure rate. Small correlations are not evidence of weak predictors; "
         "they are a class-imbalance artifact. The Random Forest captures "
         "non-linear interactions invisible to pairwise correlation.",
         cx + 0.15, cy + 12.60, COL_W - 0.3, 1.8,
         font_size=9, color=C_NAVY)

# EDA methodology
section_header(slide, "EDA METHODOLOGY", cx, cy + 14.60, COL_W, C_ACCENT)
method_text = """\
Four-step EDA using SQLite in-memory queries:

Step 1 — Data Profiling
Null audit, value-count review, zero-variance flagging. \
Confirmed days_in_fc and weather_risk have no signal.

Step 2 — Class Imbalance Assessment
Visualized 140:1 ratio; established recall as primary metric.

Step 3 — Bivariate Analysis
Failure rates by carrier, shift, route distance, and \
package type using GROUP BY SQL queries.

Step 4 — Correlation & Feature Ranking
Pearson heatmap + conditional failure-rate ranking as \
a model-free feature importance preview."""

add_text(slide, method_text, cx + 0.15, cy + 15.25, COL_W - 0.3, 9.0,
         font_size=9.5, color=C_NAVY)

# Data quality notes
section_header(slide, "DATA QUALITY DECISIONS", cx, cy + 24.50, COL_W, C_WARN)
quality_text = """\
Zero-variance exclusions: weather_risk and days_in_fc \
removed before modeling — both have identical values \
across all 3,559 rows (dataset limitation for July 2018 LA).

Data leakage prevention: damaged_on_arrival IS the target \
variable (DELIVERY_ATTEMPTED flag from scan_status). \
Using it as a feature would be circular — excluded entirely.

Feature encoding: carrier, shift, package_type label-encoded. \
Distance bucketed into operationally meaningful ranges \
(< 40 km / 40–60 km / > 60 km) before one-hot encoding."""

add_text(slide, quality_text, cx + 0.15, cy + 25.10, COL_W - 0.3, 5.8,
         font_size=9.5, color=C_NAVY)


# ═══════════════════════════════════════════════════
# COLUMN 3 — KEY FINDINGS
# ═══════════════════════════════════════════════════
cx = col_x[2]
cy = COL_Y

add_rect(slide, cx, cy, COL_W, COL_H, C_LIGHT)
section_header(slide, "KEY FINDINGS", cx, cy + 0.15, COL_W, C_DANGER)

# Chart 1: Carrier
add_text(slide, "Failure Rate by Carrier", cx + 0.15, cy + 0.90, COL_W - 0.3, 0.4,
         font_size=11, bold=True, color=C_NAVY)
add_image(slide, FIGURES_DIR / "failure_rate_by_carrier.png",
          cx + 0.10, cy + 1.30, COL_W - 0.2, 4.2)
add_text(slide,
         "carrier_D has the highest failure rate (1.39%) — nearly double the "
         "0.70% overall average. carrier_B has zero failures across 412 packages. "
         "Carrier identity is the strongest structural predictor in the dataset.",
         cx + 0.15, cy + 5.60, COL_W - 0.3, 1.3,
         font_size=9, color=C_NAVY)

# Chart 2: Shift
add_text(slide, "Failure Rate by Shift", cx + 0.15, cy + 7.10, COL_W - 0.3, 0.4,
         font_size=11, bold=True, color=C_NAVY)
add_image(slide, FIGURES_DIR / "failure_rate_by_shift.png",
          cx + 0.10, cy + 7.50, COL_W - 0.2, 4.2)
add_text(slide,
         "Morning shift drives 1.37% failure rate vs 0.55% for afternoon. "
         "In dense LA neighborhoods, morning deliveries face lower residential "
         "occupancy (commuters at work) and more access barriers. "
         "Only 2 shifts exist in this dataset — no night routes operate on these 15 LA routes.",
         cx + 0.15, cy + 11.80, COL_W - 0.3, 1.5,
         font_size=9, color=C_NAVY)

# Chart 3: Distance
add_text(slide, "Failure Rate by Route Distance", cx + 0.15, cy + 13.55, COL_W - 0.3, 0.4,
         font_size=11, bold=True, color=C_NAVY)
add_image(slide, FIGURES_DIR / "failure_rate_by_distance.png",
          cx + 0.10, cy + 13.95, COL_W - 0.2, 4.5)

# Urban Density Paradox callout
add_rect(slide, cx + 0.1, cy + 18.65, COL_W - 0.2, 2.6, RGBColor(0xFF, 0xF0, 0xEE),
         line_color=C_DANGER)
add_text(slide,
         "URBAN DENSITY PARADOX\n"
         "Routes under 40 km fail at 1.89% — the highest of any bucket. "
         "Routes over 60 km have ZERO failures. "
         "Dense LA neighborhoods present access barriers that longer suburban routes avoid: "
         "locked apartment lobbies, key-fob gates, congested Amazon Lockers, "
         "and multi-tenant buildings with no reception.",
         cx + 0.2, cy + 18.75, COL_W - 0.4, 2.4,
         font_size=9, bold=False, color=C_DANGER)

# Feature importance chart
add_text(slide, "Feature Importance Preview", cx + 0.15, cy + 21.55, COL_W - 0.3, 0.4,
         font_size=11, bold=True, color=C_NAVY)
add_image(slide, FIGURES_DIR / "feature_importance_preview.png",
          cx + 0.10, cy + 21.95, COL_W - 0.2, 5.5)
add_text(slide,
         "Model-free conditional failure-rate ranking. carrier_D and morning shift "
         "emerge as the top predictors. locker_issue shows an elevated signal despite "
         "low frequency. No single feature dominates — the Random Forest will exploit "
         "interactions between carrier, shift, and distance simultaneously.",
         cx + 0.15, cy + 27.55, COL_W - 0.3, 1.8,
         font_size=9, color=C_NAVY)

# Top findings summary
section_header(slide, "TOP 3 FINDINGS AT A GLANCE", cx, cy + 29.55, COL_W, C_DANGER)
top3 = """\
1. carrier_D fails at 2× the average rate (1.39% vs 0.70%) — structural carrier risk
2. Morning shift has 2.5× higher failure rate than afternoon (1.37% vs 0.55%)
3. Short routes fail MOST — < 40 km at 1.89% vs > 60 km at 0.00% (urban density paradox)"""
add_text(slide, top3, cx + 0.15, cy + 30.20, COL_W - 0.3, 1.0,
         font_size=9.5, bold=False, color=C_NAVY)


# ═══════════════════════════════════════════════════
# COLUMN 4 — CONCLUSIONS & NEXT STEPS
# ═══════════════════════════════════════════════════
cx = col_x[3]
cy = COL_Y

add_rect(slide, cx, cy, COL_W, COL_H, C_LIGHT)
section_header(slide, "CONCLUSIONS & NEXT STEPS", cx, cy + 0.15, COL_W, C_SAFE)

conclusions_text = """\
What This Analysis Proved

Three findings stand out from the real LMRC data — and one was a genuine \
surprise. carrier_D's elevated failure rate and morning shift's access \
barriers are expected; the counterintuitive result was route distance. \
The assumption that longer, more complex routes would fail more often is \
exactly backwards: the data shows the shortest routes (dense urban LA) \
fail almost 2× more than routes twice their length.

This matters for dispatch planning. Flagging carrier_D morning routes \
through dense neighborhoods as high-risk before the truck leaves is \
operationally tractable — the data signals are available at dispatch time.

Actionable Recommendations

1. Flag carrier_D × morning × < 40 km combinations for pre-dispatch \
supervisor review. This combination concentrates all three top risk factors.

2. Redistribute carrier_D's densest routes (< 30 km, morning) to \
carrier_B or carrier_A, which demonstrate significantly lower failure rates \
on equivalent geographic footprints.

3. Implement a morning-shift access-verification protocol: confirm delivery \
access codes, intercom numbers, or locker availability before dispatch \
for packages on sub-40 km urban routes.
"""
add_text(slide, conclusions_text, cx + 0.15, cy + 0.90, COL_W - 0.3, 9.5,
         font_size=9.5, color=C_NAVY)

# Model implications
section_header(slide, "MODEL IMPLICATIONS", cx, cy + 10.65, COL_W, C_ACCENT)
model_text = """\
Recall > Accuracy (always)
At 0.7% failure rate, accuracy is irrelevant. A model predicting \
"delivered" 100% of the time scores 99.3% accuracy with zero business value. \
Recall — what fraction of real failures does the model catch? — is \
the metric that maps to operational savings.

Why class_weight='balanced'
The RandomForest's native class_weight='balanced' reweights each failure \
by a factor of ~140× relative to a successful delivery. This prevents the \
model from defaulting to the majority class.

Why SMOTE (comparison)
Synthetic Minority Oversampling (SMOTE) generates synthetic failure samples \
in feature space. The notebook compares balanced-weight RandomForest vs \
SMOTE+RandomForest and reports recall, precision, and AUC-ROC.

Threshold Optimization
Default probability threshold (0.5) is suboptimal for recall. \
The notebook sweeps thresholds from 0.1 to 0.9 and selects the cutoff \
that maximizes recall while keeping precision above a business-defined floor.

Limitations
• 3,559 packages from 15 routes is a small training set for a 0.7% event
• Only July 2018 LA — seasonal and geographic generalization untested
• Feature set is dispatch-time only; real-time GPS signals not available"""

add_text(slide, model_text, cx + 0.15, cy + 11.30, COL_W - 0.3, 10.5,
         font_size=9.5, color=C_NAVY)

# Future work
section_header(slide, "FUTURE WORK", cx, cy + 22.10, COL_W, C_ACCENT)
future_text = """\
Scale to full LMRC dataset
The full Amazon LMRC 2021 dataset covers ~6,000 routes. Scaling from 15 \
routes to the full dataset would provide sufficient failure examples to \
train a more robust model without SMOTE dependence.

Barcelona Open Data Validation Layer
A planned extension would test whether urban-density failure patterns \
generalize to other markets using Barcelona Open Data (accident risk \
by neighborhood, traffic congestion by shift). This would validate \
whether the U.S.-specific findings translate to European delivery networks.

CrewAI Agent Integration
The agents_crew.py module is designed to generate per-package executive \
risk reports via CrewAI. The next step is integrating the trained model \
output as tool input so the agent can produce: "Package X is high-risk \
because carrier_D + morning + < 40 km. Recommend: hold for review."

Dashboard Deployment
The Streamlit operations dashboard (dashboard.py) provides real-time \
package scoring for dispatch supervisors. Deploy as a daily S&OP tool \
with route-level morning briefings."""

add_text(slide, future_text, cx + 0.15, cy + 22.75, COL_W - 0.3, 7.5,
         font_size=9.5, color=C_NAVY)

# CrewAI example output
section_header(slide, "CREWAI AGENT OUTPUT EXAMPLE", cx, cy + 30.55, COL_W,
               RGBColor(0x8E, 0x44, 0xAD))
crew_text = """\
"Package PKG_abc123 — RISK: HIGH
Carrier D, morning shift, route 31.2 km.
All three top risk factors active simultaneously.
Recommendation: Hold for pre-dispatch review.
Estimated failure cost if not intervened: $17." """
add_text(slide, crew_text, cx + 0.15, cy + 31.20, COL_W - 0.3, 1.0,
         font_size=8.5, color=C_NAVY)


# ══════════════════════════════════════════════════════════════════════════════
# FOOTER
# ══════════════════════════════════════════════════════════════════════════════
footer_y = H_IN - 0.35
add_rect(slide, 0, footer_y, W_IN, 0.35, C_NAVY)
add_text(slide,
         "Correlation One Data Analytics (DANA) — Week 12 Final Portfolio Project  |  "
         "Amazon LMRC 2021 Dataset (Public)  |  Los Angeles, July 2018  |  April 2026",
         0.3, footer_y, W_IN - 0.6, 0.30,
         font_size=8, color=C_WHITE, align=PP_ALIGN.CENTER)


# ── Save ───────────────────────────────────────────────────────────────────────
prs.save(str(OUT_PATH))
print(f"\n  ✓ Datafolio saved → {OUT_PATH}\n")
